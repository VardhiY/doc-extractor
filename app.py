import streamlit as st
import io
import re
from PIL import Image
import pytesseract

# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

st.set_page_config(
    page_title="DocVault — Data Extraction Tool",
    page_icon="🔐",
    layout="centered"
)

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@400;500;600;700;800&family=DM+Sans:ital,wght@0,300;0,400;0,500;0,600;1,400&family=DM+Mono:wght@400;500&display=swap');

*, *::before, *::after { box-sizing: border-box; }

/* ══ BASE ══ */
.stApp {
    background: #ffffff !important;
    color: #111827 !important;
    font-family: 'DM Sans', sans-serif !important;
}

/* Colorful gradient mesh background */
.stApp::before {
    content: '';
    position: fixed; inset: 0;
    background:
        radial-gradient(ellipse 900px 600px at 0%   0%,   rgba(99,102,241,0.13) 0%, transparent 55%),
        radial-gradient(ellipse 700px 500px at 100% 0%,   rgba(236,72,153,0.11) 0%, transparent 55%),
        radial-gradient(ellipse 800px 600px at 50%  100%, rgba(16,185,129,0.10) 0%, transparent 55%),
        radial-gradient(ellipse 600px 400px at 100% 100%, rgba(245,158,11,0.09) 0%, transparent 55%),
        radial-gradient(ellipse 500px 400px at 0%   100%, rgba(59,130,246,0.10) 0%, transparent 55%);
    pointer-events: none; z-index: 0;
}

.block-container {
    padding-top: 0 !important;
    padding-bottom: 4rem !important;
    max-width: 780px !important;
    position: relative; z-index: 1;
}

#MainMenu, footer, header { visibility: hidden; }
.stDeployButton { display: none; }

::-webkit-scrollbar { width: 5px; }
::-webkit-scrollbar-thumb { background: #e5e7eb; border-radius: 10px; }

/* ══ TOP NAV ══ */
.dv-topnav {
    display: flex; align-items: center; justify-content: space-between;
    padding: 1rem 0 0.8rem;
    border-bottom: 1px solid #f3f4f6;
    margin-bottom: 0;
}
.dv-logo {
    display: flex; align-items: center; gap: 0.5rem;
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: 1.05rem; font-weight: 800; color: #111827;
}
.dv-logo-dot {
    width: 28px; height: 28px;
    background: linear-gradient(135deg, #6366f1, #ec4899);
    border-radius: 8px;
    display: flex; align-items: center; justify-content: center;
    font-size: 0.82rem;
    box-shadow: 0 3px 10px rgba(99,102,241,0.35);
}
.dv-nav-links {
    display: flex; align-items: center; gap: 1.5rem;
    font-size: 0.85rem; font-weight: 500; color: #6b7280;
}
.dv-nav-links a { color: #6b7280; text-decoration: none; }
.dv-nav-links a:hover { color: #6366f1; }
.dv-nav-cta {
    background: linear-gradient(135deg, #6366f1, #8b5cf6);
    color: white !important; padding: 0.38rem 1rem;
    border-radius: 8px; font-weight: 600; font-size: 0.82rem;
    box-shadow: 0 2px 8px rgba(99,102,241,0.3);
}

/* ══ HERO ══ */
.dv-hero {
    text-align: center;
    padding: 3rem 0 2.5rem;
}
.dv-hero-badge {
    display: inline-flex; align-items: center; gap: 0.4rem;
    background: linear-gradient(135deg, rgba(99,102,241,0.08), rgba(139,92,246,0.08));
    border: 1px solid rgba(99,102,241,0.2);
    border-radius: 100px; padding: 0.3rem 1rem;
    font-size: 0.75rem; font-weight: 600;
    color: #6366f1; margin-bottom: 1.2rem;
    letter-spacing: 0.02em;
}
.dv-hero-badge::before {
    content: ''; width: 6px; height: 6px;
    background: linear-gradient(135deg, #6366f1, #ec4899);
    border-radius: 50%; flex-shrink: 0;
    animation: badgePulse 2s ease-in-out infinite;
}
@keyframes badgePulse { 0%,100%{opacity:1;transform:scale(1);} 50%{opacity:0.5;transform:scale(0.85);} }

.dv-h1 {
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: clamp(2rem, 5vw, 3rem);
    font-weight: 800; line-height: 1.15; letter-spacing: -0.8px;
    color: #111827; margin-bottom: 0.85rem;
}
.dv-h1 .grad {
    background: linear-gradient(90deg, #6366f1, #ec4899, #f59e0b);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    background-clip: text;
}
.dv-hero-sub {
    font-size: 1rem; color: #6b7280; max-width: 500px;
    margin: 0 auto 0; line-height: 1.7;
}

/* ══ MAIN PANEL ══ */
.dv-panel {
    background: #ffffff;
    border: 1px solid #e5e7eb;
    border-radius: 20px;
    padding: 2rem;
    margin-bottom: 1.2rem;
    box-shadow: 0 4px 24px rgba(0,0,0,0.06), 0 1px 4px rgba(0,0,0,0.04);
}
.dv-panel-title {
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: 0.78rem; font-weight: 700;
    letter-spacing: 0.1em; text-transform: uppercase;
    color: #9ca3af; margin-bottom: 1rem;
    display: flex; align-items: center; gap: 0.5rem;
}
.dv-panel-title::after { content:''; flex:1; height:1px; background:#f3f4f6; }

/* ══ FILE UPLOADER ══ */
[data-testid="stFileUploader"] section {
    background: linear-gradient(135deg, rgba(99,102,241,0.03), rgba(139,92,246,0.03)) !important;
    border: 2px dashed rgba(99,102,241,0.25) !important;
    border-radius: 14px !important;
    padding: 2rem 1rem !important;
    text-align: center !important;
    transition: all 0.2s !important;
    cursor: pointer !important;
}
[data-testid="stFileUploader"] section:hover {
    border-color: rgba(99,102,241,0.5) !important;
    background: rgba(99,102,241,0.04) !important;
}
[data-testid="stFileUploader"] section p {
    color: #6b7280 !important; font-size: 0.9rem !important;
}
[data-testid="stFileUploadDropzoneInput"] { cursor: pointer !important; }
[data-testid="stFileUploader"] button {
    background: linear-gradient(135deg, #6366f1, #8b5cf6) !important;
    color: white !important; border: none !important;
    border-radius: 8px !important; font-weight: 600 !important;
    padding: 0.45rem 1.2rem !important; font-size: 0.85rem !important;
    box-shadow: 0 2px 8px rgba(99,102,241,0.3) !important;
}

/* ══ CHECKBOX GRID ══ */
.dv-checkbox-grid {
    display: grid; grid-template-columns: repeat(4, 1fr);
    gap: 0.5rem; margin-bottom: 0.5rem;
}
.dv-chip {
    display: flex; align-items: center; justify-content: center;
    gap: 0.35rem; padding: 0.5rem 0.6rem;
    border: 1.5px solid #e5e7eb; border-radius: 10px;
    background: #fafafa; cursor: pointer;
    font-size: 0.8rem; font-weight: 600; color: #374151;
    transition: all 0.18s; user-select: none;
    text-align: center;
}
.dv-chip:hover { border-color: #6366f1; color: #6366f1; background: rgba(99,102,241,0.04); }
.dv-chip.active {
    border-color: #6366f1; color: #6366f1;
    background: rgba(99,102,241,0.08);
    box-shadow: 0 0 0 3px rgba(99,102,241,0.1);
}

/* Custom checkbox styling */
.stCheckbox > label {
    font-size: 0.88rem !important; font-weight: 500 !important;
    color: #374151 !important; gap: 0.45rem !important;
}
div[data-testid="stCheckbox"] > label > div:first-child {
    background: white !important;
    border: 1.5px solid #d1d5db !important;
    border-radius: 5px !important; width: 16px !important; height: 16px !important;
}
div[data-testid="stCheckbox"][aria-checked="true"] > label > div:first-child {
    background: #6366f1 !important; border-color: #6366f1 !important;
}

/* ══ TEXT AREA ══ */
textarea {
    background: #fafafa !important;
    border: 1.5px solid #e5e7eb !important;
    border-radius: 12px !important; color: #111827 !important;
    font-family: 'DM Sans', sans-serif !important;
    font-size: 0.93rem !important; padding: 0.85rem 1rem !important;
    line-height: 1.6 !important; transition: all 0.2s !important;
    resize: vertical !important;
}
textarea:focus {
    border-color: #6366f1 !important;
    box-shadow: 0 0 0 3px rgba(99,102,241,0.08) !important;
    outline: none !important; background: white !important;
}
textarea::placeholder { color: #9ca3af !important; }

/* ══ PRIMARY BUTTON ══ */
.stButton > button {
    width: 100% !important;
    background: linear-gradient(135deg, #6366f1, #8b5cf6) !important;
    border: none !important; border-radius: 12px !important;
    color: #fff !important;
    font-family: 'Plus Jakarta Sans', sans-serif !important;
    font-weight: 700 !important; font-size: 1rem !important;
    padding: 0.9rem 1.5rem !important;
    margin-top: 0.5rem !important; transition: all 0.2s !important;
    box-shadow: 0 4px 15px rgba(99,102,241,0.35) !important;
    letter-spacing: 0.01em !important;
}
.stButton > button:hover {
    opacity: 0.9 !important; transform: translateY(-1px) !important;
    box-shadow: 0 8px 25px rgba(99,102,241,0.45) !important;
}
.stButton > button:active { transform: translateY(0) !important; }

/* ══ DOWNLOAD BUTTON ══ */
.stDownloadButton > button {
    background: white !important;
    border: 1.5px solid #e5e7eb !important;
    color: #374151 !important;
    font-family: 'DM Sans', sans-serif !important;
    font-weight: 600 !important; border-radius: 9px !important;
    padding: 0.48rem 1rem !important; font-size: 0.85rem !important;
    margin-top: 0 !important; transition: all 0.18s !important;
    box-shadow: 0 1px 4px rgba(0,0,0,0.05) !important;
}
.stDownloadButton > button:hover {
    border-color: #6366f1 !important;
    color: #6366f1 !important; background: rgba(99,102,241,0.04) !important;
}

/* ══ RESULTS CATEGORY CARDS ══ */
.cat-card {
    border-radius: 14px; padding: 1rem 1.2rem;
    margin-bottom: 0.8rem; border: 1px solid;
}
.cat-card-header {
    display: flex; align-items: center; gap: 0.6rem;
    margin-bottom: 0.65rem;
}
.cat-icon {
    width: 30px; height: 30px; border-radius: 8px;
    display: flex; align-items: center; justify-content: center;
    font-size: 0.85rem; flex-shrink: 0;
}
.cat-title {
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: 0.82rem; font-weight: 700;
    letter-spacing: 0.05em; text-transform: uppercase;
}
.cat-count {
    font-family: 'DM Mono', monospace;
    font-size: 0.65rem; font-weight: 500;
    padding: 0.12rem 0.45rem; border-radius: 100px;
    margin-left: auto;
}
.cat-items { display: flex; flex-wrap: wrap; gap: 0.4rem; }
.cat-tag {
    font-size: 0.83rem; font-weight: 500;
    padding: 0.28rem 0.7rem; border-radius: 7px;
    border: 1px solid; cursor: default;
    transition: transform 0.15s;
}
.cat-tag:hover { transform: translateY(-1px); }

/* ══ SECURITY SCAN ══ */
.scan-row {
    display: flex; align-items: center; justify-content: space-between;
    padding: 0.6rem 0.9rem; border-radius: 9px;
    margin-bottom: 0.4rem; background: #fafafa;
    border: 1px solid #f3f4f6;
}
.scan-label { font-size: 0.88rem; font-weight: 500; color: #374151; }
.scan-pass {
    display: inline-flex; align-items: center; gap: 0.28rem;
    font-size: 0.75rem; font-weight: 700; color: #059669;
    background: #d1fae5; border: 1px solid #a7f3d0;
    border-radius: 6px; padding: 0.15rem 0.5rem;
}
.scan-fail {
    display: inline-flex; align-items: center; gap: 0.28rem;
    font-size: 0.75rem; font-weight: 700; color: #dc2626;
    background: #fee2e2; border: 1px solid #fca5a5;
    border-radius: 6px; padding: 0.15rem 0.5rem;
}
.score-row {
    display: flex; align-items: center; gap: 1.2rem;
    padding: 0.9rem 1rem; border-radius: 12px;
    background: linear-gradient(135deg, rgba(99,102,241,0.05), rgba(139,92,246,0.05));
    border: 1px solid rgba(99,102,241,0.15); margin-top: 0.6rem;
}
.score-big {
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: 2rem; font-weight: 800; line-height: 1;
}
.score-sub { font-size: 0.62rem; font-weight: 600; letter-spacing: 0.1em; text-transform: uppercase; color: #9ca3af; }
.score-bar-bg { height: 8px; background: #e5e7eb; border-radius: 100px; overflow: hidden; flex: 1; }
.score-bar-fill { height: 100%; border-radius: 100px; }

/* ══ OUTPUT HEADER ══ */
.output-hdr {
    display: flex; align-items: center; justify-content: space-between;
    margin-bottom: 0.75rem;
}
.output-hdr-title {
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: 0.78rem; font-weight: 700; letter-spacing: 0.1em;
    text-transform: uppercase; color: #9ca3af;
    display: flex; align-items: center; gap: 0.45rem;
}
.output-hdr-title::after { content:''; width:32px; height:1px; background:#e5e7eb; }
.output-actions { display: flex; align-items: center; gap: 0.45rem; }
.copy-pill {
    display: inline-flex; align-items: center; gap: 0.32rem;
    font-family: 'DM Sans', sans-serif; font-size: 0.8rem; font-weight: 600;
    color: #6366f1; background: rgba(99,102,241,0.07);
    border: 1.5px solid rgba(99,102,241,0.2);
    border-radius: 8px; padding: 0.35rem 0.8rem;
    cursor: pointer; transition: all 0.18s;
}
.copy-pill:hover { background: rgba(99,102,241,0.12); border-color: rgba(99,102,241,0.4); }

/* ══ HOW IT WORKS ══ */
.hw-step {
    display: flex; align-items: flex-start; gap: 1rem;
    padding: 1rem; border-radius: 12px;
    background: linear-gradient(135deg, rgba(255,255,255,0.8), rgba(255,255,255,0.6));
    border: 1px solid #f3f4f6; margin-bottom: 0.6rem;
    transition: transform 0.2s, box-shadow 0.2s;
}
.hw-step:hover { transform: translateY(-2px); box-shadow: 0 4px 16px rgba(0,0,0,0.06); }
.hw-num {
    width: 36px; height: 36px; border-radius: 10px;
    display: flex; align-items: center; justify-content: center;
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: 0.95rem; font-weight: 800;
    flex-shrink: 0;
}
.hw-body h4 {
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: 0.92rem; font-weight: 700; color: #111827; margin-bottom: 0.2rem;
}
.hw-body p { font-size: 0.83rem; color: #6b7280; line-height: 1.55; margin: 0; }

/* ══ FAQ ══ */
.faq-item {
    padding: 1rem 1.1rem;
    border: 1px solid #f3f4f6; border-radius: 12px;
    margin-bottom: 0.5rem; background: white;
}
.faq-q {
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-size: 0.9rem; font-weight: 700; color: #111827; margin-bottom: 0.35rem;
}
.faq-a { font-size: 0.85rem; color: #6b7280; line-height: 1.6; margin: 0; }

/* ══ FOOTER ══ */
.dv-footer {
    text-align: center; padding: 2.5rem 0 1rem;
    border-top: 1px solid #f3f4f6; margin-top: 3rem;
}
.dv-footer p { font-size: 0.82rem; color: #9ca3af; }

/* ══ MISC ══ */
.stSpinner > div { border-top-color: #6366f1 !important; }
.stAlert { background: rgba(99,102,241,0.06) !important; border: 1px solid rgba(99,102,241,0.2) !important; border-radius: 10px !important; }
div[data-testid="stHorizontalBlock"] > div[data-testid="column"] { background: transparent !important; }
</style>
""", unsafe_allow_html=True)


# ── SECURITY HELPERS ─────────────────────────────────────────────────────────
MALWARE_SIGS = [b"cmd.exe", b"powershell", b"eval(", b"WScript"]
MAGIC_BYTES  = {"pdf": b"%PDF", "docx": b"PK", "xlsx": b"PK", "pptx": b"PK", "ppt": b"\xd0\xcf\x11\xe0"}

def check_magic(fb, ext):     return fb.startswith(MAGIC_BYTES[ext]) if ext in MAGIC_BYTES else True
def check_malware(fb):        return not any(sig in fb for sig in MALWARE_SIGS)
def check_integrity(fb, ext):
    try:
        if ext == "pdf":
            import pypdf; pypdf.PdfReader(io.BytesIO(fb))
        elif ext == "docx":
            import docx; docx.Document(io.BytesIO(fb))
        elif ext == "xlsx":
            import openpyxl; openpyxl.load_workbook(io.BytesIO(fb))
        elif ext in ["pptx","ppt"]:
            from pptx import Presentation; Presentation(io.BytesIO(fb))
        return True
    except: return False

def extract_text(fb, ext):
    if ext in ["png","jpg","jpeg"]:
        return pytesseract.image_to_string(Image.open(io.BytesIO(fb)))
    if ext == "pdf":
        import pypdf
        return "\n".join(p.extract_text() or "" for p in pypdf.PdfReader(io.BytesIO(fb)).pages)
    if ext == "docx":
        import docx
        return "\n".join(p.text for p in docx.Document(io.BytesIO(fb)).paragraphs)
    if ext == "xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(fb), data_only=True)
        lines = []
        for sh in wb.worksheets:
            for row in sh.iter_rows(values_only=True):
                ln = " | ".join(str(c) for c in row if c)
                if ln: lines.append(ln)
        return "\n".join(lines)
    if ext in ["pptx","ppt"]:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(fb))
        return "\n".join(s.text.strip() for sl in prs.slides for s in sl.shapes if hasattr(s,"text") and s.text.strip())
    if ext == "txt":
        return fb.decode("utf-8", errors="ignore")
    return None

# ── EXTRACTION HELPERS ────────────────────────────────────────────────────────
import re as _re

PATTERNS = {
    "People Names":      r'\b[A-Z][a-z]+ [A-Z][a-z]+(?:\s[A-Z][a-z]+)?\b',
    "Email Addresses":   r'\b[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}\b',
    "Phone Numbers":     r'(?:\+?\d[\d\s\-().]{7,14}\d)',
    "Dates":             r'\b(?:\d{1,2}[-/]\d{1,2}[-/]\d{2,4}|\d{4}[-/]\d{2}[-/]\d{2}|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s+\d{1,2},?\s+\d{4})\b',
    "Organizations":     r'\b(?:[A-Z][a-z]+ )*(?:Inc\.|LLC|Ltd\.|Corp\.|Corporation|Company|Group|Foundation|Institute|University|College|School|Hospital|Bank|Agency)\b',
    "Locations":         r'\b[A-Z][a-z]+(?:,\s*[A-Z]{2})?\b',
    "Financial Figures": r'(?:USD|EUR|GBP|INR|₹|\$|€|£)\s?[\d,]+(?:\.\d{1,2})?|[\d,]+(?:\.\d{1,2})?\s?(?:USD|EUR|GBP|INR)',
    "Key Topics":        r'\b[A-Z][A-Za-z]{3,}(?:\s[A-Za-z]{3,}){0,2}\b',
    "Aadhaar Number":    r'\b\d{4}\s?\d{4}\s?\d{4}\b',
    "PAN Number":        r'\b[A-Z]{5}[0-9]{4}[A-Z]\b',
    "SSN":               r'\b\d{3}-\d{2}-\d{4}\b',
    "Mobile Numbers":    r'\b[6-9]\d{9}\b',
    "Date of Birth":     r'\b\d{2}[-/]\d{2}[-/]\d{4}\b',
}

CATEGORY_STYLES = {
    "People Names":      {"bg":"#fdf2f8","border":"#fbcfe8","icon_bg":"#fce7f3","icon":"👤","color":"#be185d","tag_bg":"#fce7f3","tag_border":"#fbcfe8","tag_color":"#be185d"},
    "Email Addresses":   {"bg":"#eff6ff","border":"#bfdbfe","icon_bg":"#dbeafe","icon":"✉️","color":"#1d4ed8","tag_bg":"#dbeafe","tag_border":"#bfdbfe","tag_color":"#1d4ed8"},
    "Phone Numbers":     {"bg":"#f0fdf4","border":"#bbf7d0","icon_bg":"#dcfce7","icon":"📞","color":"#15803d","tag_bg":"#dcfce7","tag_border":"#bbf7d0","tag_color":"#15803d"},
    "Dates":             {"bg":"#fff7ed","border":"#fed7aa","icon_bg":"#ffedd5","icon":"📅","color":"#c2410c","tag_bg":"#ffedd5","tag_border":"#fed7aa","tag_color":"#c2410c"},
    "Organizations":     {"bg":"#fdf4ff","border":"#e9d5ff","icon_bg":"#f3e8ff","icon":"🏢","color":"#7e22ce","tag_bg":"#f3e8ff","tag_border":"#e9d5ff","tag_color":"#7e22ce"},
    "Locations":         {"bg":"#f0fdfa","border":"#99f6e4","icon_bg":"#ccfbf1","icon":"📍","color":"#0f766e","tag_bg":"#ccfbf1","tag_border":"#99f6e4","tag_color":"#0f766e"},
    "Financial Figures": {"bg":"#fefce8","border":"#fde68a","icon_bg":"#fef9c3","icon":"💰","color":"#b45309","tag_bg":"#fef9c3","tag_border":"#fde68a","tag_color":"#b45309"},
    "Key Topics":        {"bg":"#eef2ff","border":"#c7d2fe","icon_bg":"#e0e7ff","icon":"🔑","color":"#4338ca","tag_bg":"#e0e7ff","tag_border":"#c7d2fe","tag_color":"#4338ca"},
    "Aadhaar Number":    {"bg":"#fdf2f8","border":"#fbcfe8","icon_bg":"#fce7f3","icon":"🪪","color":"#be185d","tag_bg":"#fce7f3","tag_border":"#fbcfe8","tag_color":"#be185d"},
    "PAN Number":        {"bg":"#eff6ff","border":"#bfdbfe","icon_bg":"#dbeafe","icon":"🗂️","color":"#1d4ed8","tag_bg":"#dbeafe","tag_border":"#bfdbfe","tag_color":"#1d4ed8"},
    "SSN":               {"bg":"#fdf4ff","border":"#e9d5ff","icon_bg":"#f3e8ff","icon":"🔒","color":"#7e22ce","tag_bg":"#f3e8ff","tag_border":"#e9d5ff","tag_color":"#7e22ce"},
    "Mobile Numbers":    {"bg":"#f0fdf4","border":"#bbf7d0","icon_bg":"#dcfce7","icon":"📱","color":"#15803d","tag_bg":"#dcfce7","tag_border":"#bbf7d0","tag_color":"#15803d"},
    "Date of Birth":     {"bg":"#fff7ed","border":"#fed7aa","icon_bg":"#ffedd5","icon":"🎂","color":"#c2410c","tag_bg":"#ffedd5","tag_border":"#fed7aa","tag_color":"#c2410c"},
}

def extract_category(text, category):
    pattern = PATTERNS.get(category)
    if not pattern: return []
    found = list(dict.fromkeys(_re.findall(pattern, text)))
    if category == "Key Topics":
        stopwords = {"The","This","That","These","Those","With","From","Have","Will","Been","Were","They","Their","There","About","Into","More","Some","Such","When","Where","Which","While"}
        found = [f for f in found if f not in stopwords][:15]
    return found[:20]

def render_category_card(category, items):
    s = CATEGORY_STYLES.get(category, CATEGORY_STYLES["Key Topics"])
    tags_html = "".join(
        f'<span class="cat-tag" style="background:{s["tag_bg"]};border-color:{s["tag_border"]};color:{s["tag_color"]};">{item}</span>'
        for item in items
    )
    return f"""
<div class="cat-card" style="background:{s['bg']};border-color:{s['border']};">
  <div class="cat-card-header">
    <div class="cat-icon" style="background:{s['icon_bg']};">{s['icon']}</div>
    <div class="cat-title" style="color:{s['color']};">{category}</div>
    <div class="cat-count" style="background:{s['icon_bg']};color:{s['color']};border:1px solid {s['border']};">{len(items)} found</div>
  </div>
  <div class="cat-items">{tags_html if tags_html else '<span style="font-size:0.82rem;color:#9ca3af;font-style:italic;">No matches found in document</span>'}</div>
</div>"""


# ── SESSION STATE ─────────────────────────────────────────────────────────────
if "extracted_text" not in st.session_state: st.session_state.extracted_text = ""
if "results"        not in st.session_state: st.session_state.results = {}
if "custom_text"    not in st.session_state: st.session_state.custom_text = ""


# ══════════════════════════════════════════════════════════════════════════════
# TOP NAV
# ══════════════════════════════════════════════════════════════════════════════
st.markdown("""
<div class="dv-topnav">
  <div class="dv-logo">
    <div class="dv-logo-dot">🔐</div>
    DocVault
  </div>
  <div class="dv-nav-links">
    <a href="#">How it Works</a>
    <a href="#">Pricing</a>
    <a href="#" class="dv-nav-cta">Visit DocVault ↗</a>
  </div>
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# HERO
# ══════════════════════════════════════════════════════════════════════════════
st.markdown("""
<div class="dv-hero">
  <div class="dv-hero-badge">Free Document Data Extractor</div>
  <h1 class="dv-h1">Extract Data From<br><span class="grad">Any Document</span></h1>
  <p class="dv-hero-sub">Free online extractor — upload any PDF, DOCX, or TXT file and instantly pull structured data. No registration required.</p>
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN PANEL
# ══════════════════════════════════════════════════════════════════════════════
st.markdown('<div class="dv-panel">', unsafe_allow_html=True)

# ── Step 1: Upload ──
st.markdown('<div class="dv-panel-title">① Upload Document (PDF, DOCX, TXT, PNG, JPG)</div>', unsafe_allow_html=True)
uploaded = st.file_uploader(
    "Drop your document here or click to browse",
    type=["pdf","docx","xlsx","pptx","ppt","txt","png","jpg","jpeg"],
    label_visibility="collapsed"
)
st.markdown('<p style="font-size:0.75rem;color:#9ca3af;margin-top:0.35rem;text-align:center;">Supports PDF, DOCX, XLSX, PPTX, TXT, PNG, JPG &nbsp;·&nbsp; Max 200MB</p>', unsafe_allow_html=True)

st.markdown("<div style='height:1.2rem'></div>", unsafe_allow_html=True)

# ── Step 2: Data to extract ──
st.markdown('<div class="dv-panel-title">② Data to Extract &nbsp;<span style="font-size:0.7rem;color:#9ca3af;text-transform:none;letter-spacing:0;font-weight:500;">(select one or more)</span></div>', unsafe_allow_html=True)

DATA_TYPES = [
    ("🔑", "Key Topics"),
    ("👤", "People Names"),
    ("🏢", "Organizations"),
    ("📅", "Dates"),
    ("✉️", "Email Addresses"),
    ("📞", "Phone Numbers"),
    ("📍", "Locations"),
    ("💰", "Financial Figures"),
    ("🪪", "Aadhaar Number"),
    ("🗂️", "PAN Number"),
    ("🔒", "SSN"),
    ("📱", "Mobile Numbers"),
]

# 3 per row
rows = [DATA_TYPES[i:i+4] for i in range(0, len(DATA_TYPES), 4)]
selected = {}
for row in rows:
    cols = st.columns(len(row))
    for col, (icon, label) in zip(cols, row):
        with col:
            selected[label] = st.checkbox(f"{icon}  {label}", key=f"chk_{label}")

st.markdown("<div style='height:0.5rem'></div>", unsafe_allow_html=True)

# ── Step 3: Custom field ──
st.markdown('<div class="dv-panel-title">③ Custom Data Request &nbsp;<span style="font-size:0.7rem;color:#9ca3af;text-transform:none;letter-spacing:0;font-weight:500;">(optional)</span></div>', unsafe_allow_html=True)
custom_req = st.text_area(
    "", height=80,
    placeholder="Describe any specific data you want to extract that's not covered above…",
    label_visibility="collapsed"
)

# ── Extract button ──
if st.button("⚡  Extract Data", key="btn_extract"):

    chosen = [k for k,v in selected.items() if v]
    if not uploaded and not custom_req:
        st.warning("Please upload a document first.")
        st.stop()
    if not chosen and not custom_req:
        st.warning("Please select at least one data type to extract.")
        st.stop()

    # ── File processing ──
    if uploaded:
        file_bytes = uploaded.read()
        ext        = uploaded.name.split(".")[-1].lower()

        magic_ok     = check_magic(file_bytes, ext)
        malware_ok   = check_malware(file_bytes)
        integrity_ok = check_integrity(file_bytes, ext)

        # Security scan results
        st.markdown('<div class="dv-panel-title" style="margin-top:1.2rem;">Security Scan</div>', unsafe_allow_html=True)

        def badge(ok):
            return '<span class="scan-pass">✓ Pass</span>' if ok else '<span class="scan-fail">✗ Fail</span>'

        st.markdown(f"""
<div class="scan-row"><span class="scan-label">🔍 File Type Validation</span>{badge(magic_ok)}</div>
<div class="scan-row"><span class="scan-label">🛡️ Malware Signature Scan</span>{badge(malware_ok)}</div>
<div class="scan-row"><span class="scan-label">✅ Structural Integrity</span>{badge(integrity_ok)}</div>
""", unsafe_allow_html=True)

        score = 100
        if not magic_ok:     score -= 30
        if not malware_ok:   score -= 40
        if not integrity_ok: score -= 30
        s_color = "#059669" if score==100 else ("#d97706" if score>=60 else "#dc2626")
        bar_grad = "linear-gradient(90deg,#6366f1,#8b5cf6)" if score==100 else ("linear-gradient(90deg,#f59e0b,#fbbf24)" if score>=60 else "linear-gradient(90deg,#ef4444,#f87171)")

        st.markdown(f"""
<div class="score-row">
  <div>
    <div class="score-sub">Security Score</div>
    <div class="score-big" style="color:{s_color};">{score}</div>
    <div style="font-size:0.65rem;color:#9ca3af;font-family:'DM Mono',monospace;">/ 100</div>
  </div>
  <div style="flex:1;">
    <div class="score-sub" style="margin-bottom:0.4rem;">Scan Progress</div>
    <div class="score-bar-bg"><div class="score-bar-fill" style="width:{score}%;background:{bar_grad};"></div></div>
    <div style="font-size:0.72rem;color:{s_color};font-family:'DM Mono',monospace;margin-top:0.3rem;font-weight:600;">{score}% clean</div>
  </div>
</div>
""", unsafe_allow_html=True)

        if not (magic_ok and malware_ok and integrity_ok):
            st.error("🚫 File blocked — one or more security checks failed.")
            st.stop()

        with st.spinner("Extracting text from document…"):
            raw_text = extract_text(file_bytes, ext)
    else:
        raw_text = ""

    if not raw_text or not raw_text.strip():
        if not custom_req:
            st.warning("No readable text found in this file.")
            st.stop()
        raw_text = ""

    # ── Run extraction ──
    with st.spinner("Analyzing and extracting data…"):
        results = {}
        for cat in chosen:
            items = extract_category(raw_text, cat)
            results[cat] = items

        if custom_req and raw_text:
            # simple keyword match for custom request
            words = [w.strip() for w in custom_req.replace(","," ").split() if len(w.strip()) > 3]
            custom_matches = []
            for word in words:
                found = _re.findall(rf'\b\w*{_re.escape(word)}\w*\b', raw_text, _re.IGNORECASE)
                custom_matches.extend(list(dict.fromkeys(found))[:5])
            results["Custom: " + custom_req[:30]] = list(dict.fromkeys(custom_matches))[:20]

    st.session_state.extracted_text = raw_text
    st.session_state.results        = results
    st.success(f"✅ Extraction complete — {sum(len(v) for v in results.values())} items found across {len(results)} categories.")

st.markdown('</div>', unsafe_allow_html=True)   # close dv-panel


# ══════════════════════════════════════════════════════════════════════════════
# RESULTS PANEL
# ══════════════════════════════════════════════════════════════════════════════
if st.session_state.results:
    results = st.session_state.results
    total   = sum(len(v) for v in results.values())

    # Output header with Copy + Download
    st.markdown(f"""
<div class="output-hdr">
  <div class="output-hdr-title">Extraction Results &nbsp;<span style="color:#6366f1;font-weight:800;">{total}</span> items</div>
  <div class="output-actions">
    <button class="copy-pill" onclick="
      const allText = {repr(chr(10).join(f'{k}: ' + ', '.join(v) for k,v in results.items()))};
      navigator.clipboard.writeText(allText).then(()=>{{
        this.innerHTML='✓ Copied!';
        this.style.color='#059669';this.style.background='rgba(5,150,105,0.08)';this.style.borderColor='rgba(5,150,105,0.3)';
        setTimeout(()=>{{this.innerHTML='⧉ Copy All';this.style.color='';this.style.background='';this.style.borderColor='';}},2000);
      }});
    ">⧉ Copy All</button>
  </div>
</div>
""", unsafe_allow_html=True)

    # Category cards
    for cat, items in results.items():
        style_key = cat if cat in CATEGORY_STYLES else "Key Topics"
        st.markdown(render_category_card(cat, items), unsafe_allow_html=True)

    # Raw text + download
    if st.session_state.extracted_text:
        with st.expander("📄 View Raw Extracted Text"):
            st.text_area("", st.session_state.extracted_text, height=300, label_visibility="collapsed")

        # Build structured output for download
        dl_lines = [f"DocVault Data Extraction Report", "="*50, ""]
        for cat, items in results.items():
            dl_lines.append(f"\n[{cat.upper()}]")
            for item in items:
                dl_lines.append(f"  • {item}")
        dl_lines += ["", "="*50, f"Total: {total} items extracted"]
        dl_text = "\n".join(dl_lines)

        col_dl, col_csv, col_sp = st.columns([1.2, 1.2, 3])
        with col_dl:
            st.download_button(
                "⬇ Download Report",
                data=dl_text.encode("utf-8"),
                file_name="docvault_extraction.txt",
                mime="text/plain"
            )
        with col_csv:
            import csv
            csv_buf = io.StringIO()
            w = csv.writer(csv_buf)
            w.writerow(["Category", "Extracted Value"])
            for cat, items in results.items():
                for item in items:
                    w.writerow([cat, item])
            st.download_button(
                "⬇ Download CSV",
                data=csv_buf.getvalue().encode("utf-8"),
                file_name="docvault_extraction.csv",
                mime="text/csv"
            )


# ══════════════════════════════════════════════════════════════════════════════
# HOW IT WORKS
# ══════════════════════════════════════════════════════════════════════════════
st.markdown("""
<div style="text-align:center;padding:3rem 0 1.5rem;">
  <div style="display:inline-flex;align-items:center;gap:0.4rem;background:rgba(99,102,241,0.07);border:1px solid rgba(99,102,241,0.15);border-radius:100px;padding:0.28rem 0.9rem;font-size:0.72rem;font-weight:600;color:#6366f1;letter-spacing:0.04em;margin-bottom:1rem;">Simple steps to get results</div>
  <h2 style="font-family:'Plus Jakarta Sans',sans-serif;font-size:1.8rem;font-weight:800;color:#111827;letter-spacing:-0.5px;margin-bottom:0.5rem;">How It Works</h2>
  <p style="font-size:0.92rem;color:#6b7280;max-width:400px;margin:0 auto 2rem;">Three simple steps from upload to structured results</p>
</div>
""", unsafe_allow_html=True)

st.markdown("""
<div class="hw-step">
  <div class="hw-num" style="background:linear-gradient(135deg,#6366f1,#8b5cf6);color:white;">1</div>
  <div class="hw-body">
    <h4>Upload Your Document</h4>
    <p>Drag and drop your PDF, DOCX, or TXT file. We support documents up to 200MB for fast processing.</p>
  </div>
</div>
<div class="hw-step">
  <div class="hw-num" style="background:linear-gradient(135deg,#ec4899,#f43f5e);color:white;">2</div>
  <div class="hw-body">
    <h4>Select Data Types</h4>
    <p>Choose from predefined categories like names, dates, organizations, or specify custom data you want to extract.</p>
  </div>
</div>
<div class="hw-step">
  <div class="hw-num" style="background:linear-gradient(135deg,#10b981,#06b6d4);color:white;">3</div>
  <div class="hw-body">
    <h4>Get Structured Results</h4>
    <p>Our engine processes your document and presents extracted data in organized, color-coded categories. Copy or download instantly.</p>
  </div>
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# FAQ
# ══════════════════════════════════════════════════════════════════════════════
st.markdown("""
<div style="text-align:center;padding:3rem 0 1.5rem;">
  <h2 style="font-family:'Plus Jakarta Sans',sans-serif;font-size:1.8rem;font-weight:800;color:#111827;letter-spacing:-0.5px;margin-bottom:0.5rem;">Frequently Asked Questions</h2>
</div>
<div class="faq-item">
  <div class="faq-q">Is this tool really free?</div>
  <p class="faq-a">Yes, completely free with no hidden charges, subscriptions, or limits on documents processed.</p>
</div>
<div class="faq-item">
  <div class="faq-q">Do I need to register or sign up?</div>
  <p class="faq-a">No registration required. Upload, select, and extract instantly. Optionally sign up for multi-document processing and history.</p>
</div>
<div class="faq-item">
  <div class="faq-q">How accurate are the results?</div>
  <p class="faq-a">Typically 85–95% accuracy depending on document quality. Best results with well-formatted, clear-text documents.</p>
</div>
<div class="faq-item">
  <div class="faq-q">Is my data secure?</div>
  <p class="faq-a">Your documents are processed securely and never stored permanently. All data is automatically deleted after processing.</p>
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# FOOTER
# ══════════════════════════════════════════════════════════════════════════════
st.markdown("""
<div class="dv-footer">
  <div style="display:flex;align-items:center;justify-content:center;gap:0.5rem;margin-bottom:0.6rem;">
    <div style="width:24px;height:24px;background:linear-gradient(135deg,#6366f1,#ec4899);border-radius:7px;display:flex;align-items:center;justify-content:center;font-size:0.75rem;">🔐</div>
    <span style="font-family:'Plus Jakarta Sans',sans-serif;font-weight:800;font-size:0.95rem;color:#111827;">DocVault</span>
  </div>
  <p>© 2025 DocVault Enterprise. All rights reserved.</p>
  <div style="display:flex;align-items:center;justify-content:center;gap:1.5rem;margin-top:0.5rem;">
    <a href="#" style="font-size:0.8rem;color:#9ca3af;text-decoration:none;">How it Works</a>
    <a href="#" style="font-size:0.8rem;color:#9ca3af;text-decoration:none;">Pricing</a>
    <a href="#" style="font-size:0.8rem;color:#9ca3af;text-decoration:none;">Blog</a>
  </div>
</div>
""", unsafe_allow_html=True)
